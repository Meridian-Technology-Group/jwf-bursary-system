"""
JWF Bursary System — Buyer Deck Generator
Minimalist infographic-style PPTX with Meridian Technology Group branding.

Run from the venv at /tmp/pptx-venv:
    source /tmp/pptx-venv/bin/activate && python build_deck.py
Output: jwf-bursary-proposal.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- Brand palette (extracted from meridiantech.group) ----------
NAVY        = RGBColor(0x0A, 0x1F, 0x3D)   # primary
NAVY_DARK   = RGBColor(0x04, 0x12, 0x2A)
NAVY_LIGHT  = RGBColor(0x1A, 0x38, 0x66)
ACCENT      = RGBColor(0x11, 0xB8, 0xC6)   # cyan accent
ACCENT_GLOW = RGBColor(0x6B, 0xDE, 0xEA)
TEXT_DARK   = RGBColor(0x1D, 0x29, 0x37)
TEXT_MUTED  = RGBColor(0x81, 0x8D, 0x9D)
BG_MUTED    = RGBColor(0xF2, 0xF4, 0xF7)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LINE_LIGHT  = RGBColor(0xE5, 0xE7, 0xEB)

FONT = "Inter"
FONT_FALLBACK = "Calibri"

# ---------- Presentation setup (16:9 widescreen) ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- Helpers ----------
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, italic=False,
             line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return tb


def add_accent_bar(slide, x, y, length=Inches(0.55), thickness=Pt(3.5)):
    """Short cyan bar — used as a section accent."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y,
                                 length, Emu(int(thickness * 12700)))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    return bar


def add_footer(slide, page_num, total):
    # Thin bottom rule
    add_rect(slide, Inches(0.6), Inches(7.05),
             SW - Inches(1.2), Emu(int(0.5 * 12700)), LINE_LIGHT)
    add_text(slide, Inches(0.6), Inches(7.15), Inches(6),
             Inches(0.3), "Meridian Technology Group  ·  John Whitgift Foundation",
             size=9, color=TEXT_MUTED)
    add_text(slide, SW - Inches(2), Inches(7.15), Inches(1.4),
             Inches(0.3), f"{page_num:02d} / {total:02d}",
             size=9, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


def add_logo(slide, x, y):
    """Minimalist mark: navy square with cyan dot + 'MERIDIAN' wordmark."""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y,
                                 Inches(0.32), Inches(0.32))
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 x + Inches(0.20), y + Inches(0.20),
                                 Inches(0.09), Inches(0.09))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    dot.line.fill.background()
    add_text(slide, x + Inches(0.42), y + Inches(0.02), Inches(2.2),
             Inches(0.32), "MERIDIAN", size=11, bold=True, color=NAVY)
    add_text(slide, x + Inches(0.42), y + Inches(0.18), Inches(2.2),
             Inches(0.20), "TECHNOLOGY GROUP", size=7,
             color=TEXT_MUTED)


def slide_header(slide, eyebrow, title, page_num, total):
    add_logo(slide, Inches(0.6), Inches(0.45))
    add_accent_bar(slide, Inches(0.6), Inches(1.25))
    add_text(slide, Inches(0.6), Inches(1.32), Inches(8), Inches(0.32),
             eyebrow.upper(), size=10, bold=True, color=ACCENT)
    add_text(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(0.85),
             title, size=30, bold=True, color=NAVY, line_spacing=1.05)
    add_footer(slide, page_num, total)


# ---------- Slides ----------
TOTAL = 10


# Slide 1 — Title (full-bleed navy with cyan accent)
def slide_title():
    s = prs.slides.add_slide(BLANK)
    # Background
    add_rect(s, 0, 0, SW, SH, NAVY)
    # Accent diagonal block
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                SW - Inches(0.18), 0, Inches(0.18), SH)
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()
    # Logo (white version)
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.7),
                             Inches(0.36), Inches(0.36))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.fill.background()
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             Inches(1.02), Inches(0.92),
                             Inches(0.10), Inches(0.10))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    dot.line.fill.background()
    add_text(s, Inches(1.30), Inches(0.74), Inches(4), Inches(0.32),
             "MERIDIAN", size=12, bold=True, color=WHITE)
    add_text(s, Inches(1.30), Inches(0.92), Inches(4), Inches(0.22),
             "TECHNOLOGY GROUP", size=8, color=ACCENT_GLOW)

    # Eyebrow
    add_accent_bar(s, Inches(0.8), Inches(2.6), length=Inches(0.7))
    add_text(s, Inches(0.8), Inches(2.7), Inches(8), Inches(0.4),
             "PROPOSAL  ·  APRIL 2026", size=11, bold=True, color=ACCENT)

    # Main title
    add_text(s, Inches(0.8), Inches(3.05), Inches(11.5), Inches(1.4),
             "Bursary Assessment Platform", size=54, bold=True, color=WHITE,
             line_spacing=1.0)
    add_text(s, Inches(0.8), Inches(4.15), Inches(11.5), Inches(0.8),
             "A purpose-built replacement for Symplectic Grant Tracker",
             size=22, color=ACCENT_GLOW, line_spacing=1.1)

    # Footer
    add_text(s, Inches(0.8), Inches(6.6), Inches(8), Inches(0.3),
             "Prepared for", size=9, bold=True, color=ACCENT)
    add_text(s, Inches(0.8), Inches(6.78), Inches(8), Inches(0.4),
             "John Whitgift Foundation", size=18, bold=True, color=WHITE)
    add_text(s, SW - Inches(4.5), Inches(6.78), Inches(3.7), Inches(0.4),
             "29 April 2026", size=14, color=ACCENT_GLOW,
             align=PP_ALIGN.RIGHT)


# Slide 2 — Where we are today (status infographic)
def slide_status():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Status", "Where we are today", 2, TOTAL)

    add_text(s, Inches(0.6), Inches(2.6), Inches(12), Inches(0.4),
             "The system is substantially built. Today is about hardening, contracts, and migration.",
             size=14, color=TEXT_MUTED)

    # Four phase tiles
    phases = [
        ("01", "Foundations", "Schema, auth, layout shells", True),
        ("02", "Core Platform", "Portal, assessment engine, recommendations", True),
        ("03", "Operational", "Exports, reports, audit, GDPR, PDF", True),
        ("04", "Production", "Hardening, UAT, go-live", False),
    ]

    tile_w = Inches(2.85)
    tile_h = Inches(2.6)
    gap = Inches(0.2)
    start_x = Inches(0.6)
    y = Inches(3.4)

    for i, (num, title, desc, done) in enumerate(phases):
        x = start_x + (tile_w + gap) * i
        # Tile background
        bg = NAVY if done else BG_MUTED
        add_rect(s, x, y, tile_w, tile_h, bg)
        # Top accent bar
        bar_color = ACCENT if done else TEXT_MUTED
        add_rect(s, x, y, tile_w, Inches(0.06), bar_color)
        # Phase number
        num_color = ACCENT if done else TEXT_MUTED
        add_text(s, x + Inches(0.3), y + Inches(0.3), tile_w - Inches(0.6),
                 Inches(0.6), num, size=36, bold=True, color=num_color,
                 line_spacing=1.0)
        # Title
        title_color = WHITE if done else NAVY
        add_text(s, x + Inches(0.3), y + Inches(1.05), tile_w - Inches(0.6),
                 Inches(0.4), title, size=15, bold=True, color=title_color)
        # Status pill
        status_text = "COMPLETE" if done else "IN FLIGHT"
        status_color = ACCENT if done else NAVY
        add_text(s, x + Inches(0.3), y + Inches(1.45), tile_w - Inches(0.6),
                 Inches(0.3), status_text, size=8, bold=True,
                 color=status_color)
        # Description
        desc_color = ACCENT_GLOW if done else TEXT_MUTED
        add_text(s, x + Inches(0.3), y + Inches(1.78), tile_w - Inches(0.6),
                 Inches(0.7), desc, size=11, color=desc_color,
                 line_spacing=1.25)

    # Bottom stats row
    stats_y = Inches(6.25)
    stats = [("36", "routes built"),
             ("100%", "Must-Have features"),
             ("0", "TypeScript errors"),
             ("≥10", "calc test cases")]
    sw = Inches(2.85)
    for i, (n, label) in enumerate(stats):
        x = start_x + (sw + gap) * i
        add_text(s, x, stats_y, sw, Inches(0.45), n,
                 size=22, bold=True, color=NAVY)
        add_text(s, x, stats_y + Inches(0.45), sw, Inches(0.3), label,
                 size=10, color=TEXT_MUTED)


# Slide 3 — Timeline to production (horizontal gantt-style)
def slide_timeline():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Timeline", "Path to production · end of June 2026", 3, TOTAL)

    add_text(s, Inches(0.6), Inches(2.6), Inches(12), Inches(0.4),
             "Nine-week runway · four gates",
             size=14, color=TEXT_MUTED)

    # Timeline track
    track_y = Inches(3.6)
    track_x = Inches(0.8)
    track_w = SW - Inches(1.6)
    add_rect(s, track_x, track_y, track_w, Inches(0.04), LINE_LIGHT)

    gates = [
        ("G1", "Initial Release", "Beta on staging", "11 May", 0.05),
        ("G2", "User Acceptance", "Assessor validation\n≥10 calc cases", "18 May → 12 Jun", 0.35),
        ("G3", "Pre-Prod Hardening", "GDPR · perf · security", "15 → 22 Jun", 0.70),
        ("G4", "Production Go-Live", "Cutover · 2-wk hypercare", "End June", 0.95),
    ]

    for code, title, desc, date, pos in gates:
        cx = track_x + Emu(int(int(track_w) * pos))
        # Node
        node = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  cx - Inches(0.18), track_y - Inches(0.16),
                                  Inches(0.36), Inches(0.36))
        node.fill.solid()
        node.fill.fore_color.rgb = ACCENT
        node.line.color.rgb = NAVY
        node.line.width = Pt(2)
        # Code label inside-ish (above)
        add_text(s, cx - Inches(0.5), track_y - Inches(0.7), Inches(1),
                 Inches(0.3), code, size=11, bold=True, color=ACCENT,
                 align=PP_ALIGN.CENTER)
        # Date directly above code
        add_text(s, cx - Inches(1.3), track_y - Inches(1.05), Inches(2.6),
                 Inches(0.3), date, size=10, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
        # Title below
        add_text(s, cx - Inches(1.3), track_y + Inches(0.4), Inches(2.6),
                 Inches(0.4), title, size=13, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
        # Description
        add_text(s, cx - Inches(1.3), track_y + Inches(0.85), Inches(2.6),
                 Inches(0.8), desc, size=10, color=TEXT_MUTED,
                 align=PP_ALIGN.CENTER, line_spacing=1.2)

    # Footnote
    add_rect(s, Inches(0.6), Inches(6.2), Inches(12.13), Inches(0.55), BG_MUTED)
    add_accent_bar(s, Inches(0.6), Inches(6.2),
                   length=Emu(int(0.05 * 914400)),
                   thickness=Pt(0))
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.2),
                             Inches(0.06), Inches(0.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    add_text(s, Inches(0.85), Inches(6.32), Inches(11.8), Inches(0.32),
             "Migration from Grant Tracker runs as a parallel workstream — see slide 9.",
             size=11, italic=True, color=TEXT_DARK)


# Slide 4 — Commercial headline (3 big tiles)
def slide_commercial_headline():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Commercial", "Headline terms", 4, TOTAL)

    items = [
        ("£5,000", "Initial Build", "One-off · 100% on contract signature"),
        ("£7,000", "Annual License", "Year 1 invoiced on go-live · annually thereafter"),
        ("Custom", "Migration", "Discovery-led · separately scoped"),
    ]

    tile_w = Inches(4.0)
    tile_h = Inches(3.6)
    gap = Inches(0.13)
    total_w = tile_w * 3 + gap * 2
    start_x = (SW - total_w) / 2
    y = Inches(2.85)

    for i, (figure, label, desc) in enumerate(items):
        x = start_x + (tile_w + gap) * i
        # Tile
        is_primary = i == 1
        bg = NAVY if is_primary else WHITE
        add_rect(s, x, y, tile_w, tile_h, bg,
                 line=(LINE_LIGHT if not is_primary else None))
        # Top accent bar
        add_rect(s, x, y, tile_w, Inches(0.08), ACCENT)
        # Label
        label_color = ACCENT if is_primary else ACCENT
        add_text(s, x + Inches(0.4), y + Inches(0.45), tile_w - Inches(0.8),
                 Inches(0.4), label.upper(), size=11, bold=True,
                 color=label_color)
        # Big figure
        fig_color = WHITE if is_primary else NAVY
        add_text(s, x + Inches(0.4), y + Inches(1.0), tile_w - Inches(0.8),
                 Inches(1.5), figure, size=72, bold=True, color=fig_color,
                 line_spacing=1.0)
        # Description
        desc_color = ACCENT_GLOW if is_primary else TEXT_MUTED
        add_text(s, x + Inches(0.4), y + Inches(2.7), tile_w - Inches(0.8),
                 Inches(0.7), desc, size=12, color=desc_color,
                 line_spacing=1.25)

    # Payment cadence strip
    cadence_y = Inches(6.6)
    cadence_h = Inches(0.42)
    add_rect(s, Inches(0.6), cadence_y, SW - Inches(1.2), cadence_h, BG_MUTED)
    add_rect(s, Inches(0.6), cadence_y, Inches(0.06), cadence_h, ACCENT)
    add_text(s, Inches(0.85), cadence_y + Inches(0.06), Inches(2.3),
             Inches(0.3), "PAYMENT CADENCE", size=9, bold=True, color=ACCENT)
    add_text(s, Inches(2.85), cadence_y + Inches(0.06),
             SW - Inches(3.45), Inches(0.3),
             "Build on contract signature  ·  Year 1 licence on go-live (end June 2026)  ·  Net 30 day terms",
             size=11, color=TEXT_DARK)


# Slide 5 — License fee build-up (transparent breakdown)
def slide_license_breakdown():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Licence Fee", "Transparent build-up", 5, TOTAL)

    add_text(s, Inches(0.6), Inches(2.6), Inches(12), Inches(0.4),
             "Built bottom-up from real costs · no hidden margin on infrastructure",
             size=14, color=TEXT_MUTED)

    # Left column: Infrastructure stacked bar with line items
    left_x = Inches(0.6)
    col_w = Inches(6.0)
    y0 = Inches(3.25)

    add_text(s, left_x, y0, col_w, Inches(0.4),
             "INFRASTRUCTURE", size=11, bold=True, color=ACCENT)

    add_text(s, left_x, y0 + Inches(0.35), col_w, Inches(0.4),
             "~ £1,000 / year", size=26, bold=True, color=NAVY)

    items = [
        ("Vercel Pro", "Production hosting · edge · zero-downtime", "£200"),
        ("Supabase Pro", "Postgres · auth · storage · daily backup", "£240"),
        ("Resend", "Transactional email", "£190"),
        ("Sentry", "Error monitoring", "£250"),
        ("Domain & SSL", "Custom domain", "£20"),
        ("Headroom", "Traffic & overage buffer", "£100"),
    ]

    row_y = y0 + Inches(1.05)
    for name, desc, cost in items:
        add_rect(s, left_x, row_y, col_w, Inches(0.42),
                 BG_MUTED if items.index((name, desc, cost)) % 2 == 0 else WHITE)
        # Cyan dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 left_x + Inches(0.15), row_y + Inches(0.18),
                                 Inches(0.08), Inches(0.08))
        dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        add_text(s, left_x + Inches(0.35), row_y + Inches(0.08),
                 Inches(1.6), Inches(0.3), name, size=11, bold=True,
                 color=NAVY)
        add_text(s, left_x + Inches(1.95), row_y + Inches(0.1),
                 col_w - Inches(2.7), Inches(0.3), desc, size=10,
                 color=TEXT_MUTED)
        add_text(s, left_x + col_w - Inches(0.7), row_y + Inches(0.08),
                 Inches(0.6), Inches(0.3), cost, size=11, bold=True,
                 color=NAVY, align=PP_ALIGN.RIGHT)
        row_y = row_y + Inches(0.42)

    # Right column: Maintenance + total
    right_x = Inches(7.0)
    rcol_w = Inches(5.7)

    add_text(s, right_x, y0, rcol_w, Inches(0.4),
             "MAINTENANCE", size=11, bold=True, color=ACCENT)
    add_text(s, right_x, y0 + Inches(0.35), rcol_w, Inches(0.4),
             "£6,000 / year", size=26, bold=True, color=NAVY)
    add_text(s, right_x, y0 + Inches(1.0), rcol_w, Inches(0.4),
             "£500 / month allocation", size=12, color=TEXT_MUTED)

    add_text(s, right_x, y0 + Inches(1.45), rcol_w, Inches(2.0),
             "Bug fixes\n"
             "Security patches & dependency updates\n"
             "Reference table tweaks\n"
             "Minor enhancements\n"
             "Ad-hoc support requests\n"
             "Quarterly review meetings",
             size=12, color=TEXT_DARK, line_spacing=1.5)

    # Total panel — navy
    total_y = y0 + Inches(3.0) - Inches(0.1)
    add_rect(s, right_x, total_y, rcol_w, Inches(0.95), NAVY)
    add_rect(s, right_x, total_y, Inches(0.08), Inches(0.95), ACCENT)
    add_text(s, right_x + Inches(0.3), total_y + Inches(0.12),
             rcol_w - Inches(0.6), Inches(0.3), "TOTAL ANNUAL LICENCE",
             size=10, bold=True, color=ACCENT)
    add_text(s, right_x + Inches(0.3), total_y + Inches(0.38),
             rcol_w - Inches(0.6), Inches(0.5), "£7,000 / year",
             size=28, bold=True, color=WHITE)


# Slide 6 — What maintenance covers (split: included / excluded)
def slide_maintenance_scope():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Scope", "What maintenance covers", 6, TOTAL)

    col_w = Inches(6.0)
    gap = Inches(0.13)
    start_x = (SW - col_w * 2 - gap) / 2
    y = Inches(2.85)
    h = Inches(3.95)

    # Included
    x = start_x
    add_rect(s, x, y, col_w, h, BG_MUTED)
    add_rect(s, x, y, col_w, Inches(0.08), ACCENT)
    add_text(s, x + Inches(0.4), y + Inches(0.4), col_w - Inches(0.8),
             Inches(0.35), "INCLUDED", size=11, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.4), y + Inches(0.75), col_w - Inches(0.8),
             Inches(0.5), "No additional charge", size=20, bold=True,
             color=NAVY)
    incl = [
        "Bug fixes & regressions",
        "Security patches",
        "Dependency updates",
        "Reference table tweaks",
        "Minor enhancements",
        "Ad-hoc support questions",
        "Annual GDPR review support",
    ]
    iy = y + Inches(1.55)
    for item in incl:
        # Cyan tick mark
        tick = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  x + Inches(0.45), iy + Inches(0.08),
                                  Inches(0.14), Inches(0.14))
        tick.fill.solid(); tick.fill.fore_color.rgb = ACCENT
        tick.line.fill.background()
        add_text(s, x + Inches(0.75), iy, col_w - Inches(1.2), Inches(0.32),
                 item, size=12, color=TEXT_DARK)
        iy = iy + Inches(0.32)

    # Excluded
    x = start_x + col_w + gap
    add_rect(s, x, y, col_w, h, WHITE, line=LINE_LIGHT)
    add_rect(s, x, y, col_w, Inches(0.08), TEXT_MUTED)
    add_text(s, x + Inches(0.4), y + Inches(0.4), col_w - Inches(0.8),
             Inches(0.35), "QUOTED SEPARATELY", size=11, bold=True,
             color=TEXT_MUTED)
    add_text(s, x + Inches(0.4), y + Inches(0.75), col_w - Inches(0.8),
             Inches(0.5), "Custom scope", size=20, bold=True, color=NAVY)
    excl = [
        "New major features",
        "Phase 2 roadmap items",
        "Calculation model changes",
        "Integrations with external systems",
        "Data migrations beyond Grant Tracker",
        "Bespoke reporting builds",
    ]
    iy = y + Inches(1.55)
    for item in excl:
        # Empty circle
        ring = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  x + Inches(0.45), iy + Inches(0.08),
                                  Inches(0.14), Inches(0.14))
        ring.fill.solid(); ring.fill.fore_color.rgb = WHITE
        ring.line.color.rgb = TEXT_MUTED
        ring.line.width = Pt(1.25)
        add_text(s, x + Inches(0.75), iy, col_w - Inches(1.2), Inches(0.32),
                 item, size=12, color=TEXT_DARK)
        iy = iy + Inches(0.32)


# Slide 7 — Support model
def slide_support():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Support", "Reliable communication, real coverage", 7, TOTAL)

    # Top: two highlight tiles
    tile_y = Inches(2.7)
    tile_h = Inches(1.45)
    tile_w = Inches(6.0)
    gap = Inches(0.13)
    start_x = (SW - tile_w * 2 - gap) / 2

    # Tile 1: Freshdesk
    add_rect(s, start_x, tile_y, tile_w, tile_h, NAVY)
    add_rect(s, start_x, tile_y, Inches(0.08), tile_h, ACCENT)
    add_text(s, start_x + Inches(0.35), tile_y + Inches(0.18),
             tile_w - Inches(0.6), Inches(0.3), "TICKETING",
             size=10, bold=True, color=ACCENT)
    add_text(s, start_x + Inches(0.35), tile_y + Inches(0.45),
             tile_w - Inches(0.6), Inches(0.5), "Freshdesk",
             size=24, bold=True, color=WHITE)
    add_text(s, start_x + Inches(0.35), tile_y + Inches(0.95),
             tile_w - Inches(0.6), Inches(0.5),
             "Email-to-ticket · self-service portal · full SLA audit trail",
             size=11, color=ACCENT_GLOW, line_spacing=1.3)

    # Tile 2: Team coverage
    x2 = start_x + tile_w + gap
    add_rect(s, x2, tile_y, tile_w, tile_h, BG_MUTED)
    add_rect(s, x2, tile_y, Inches(0.08), tile_h, ACCENT)
    add_text(s, x2 + Inches(0.35), tile_y + Inches(0.18),
             tile_w - Inches(0.6), Inches(0.3), "COVERAGE",
             size=10, bold=True, color=ACCENT)
    add_text(s, x2 + Inches(0.35), tile_y + Inches(0.45),
             tile_w - Inches(0.6), Inches(0.5), "Team-backed",
             size=24, bold=True, color=NAVY)
    add_text(s, x2 + Inches(0.35), tile_y + Inches(0.95),
             tile_w - Inches(0.6), Inches(0.5),
             "Brian as primary contact · outsourced engineering team scaled on demand",
             size=11, color=TEXT_MUTED, line_spacing=1.3)

    # SLA table heading
    table_y = Inches(4.45)
    add_text(s, Inches(0.6), table_y, Inches(8), Inches(0.3),
             "RESPONSE SLA",  size=10, bold=True, color=ACCENT)
    add_text(s, Inches(0.6), table_y + Inches(0.25), Inches(8), Inches(0.4),
             "Severity-based, working hours Mon–Fri 09:00–17:30 UK",
             size=12, color=TEXT_MUTED)

    # SLA rows
    sla = [
        ("CRITICAL",  "System down during a round",       "4 working hours", "1 working day",        NAVY,  WHITE),
        ("HIGH",      "Major feature broken",              "1 working day",   "3 working days",       NAVY_LIGHT, WHITE),
        ("MEDIUM",    "Minor bug · workaround exists",     "2 working days",  "Next maintenance cycle", BG_MUTED, NAVY),
        ("LOW",       "Question · suggestion",             "5 working days",  "Quarterly review",     WHITE, NAVY),
    ]
    rx = Inches(0.6)
    rw = SW - Inches(1.2)
    rh = Inches(0.36)
    rrow_y = table_y + Inches(0.85)
    # column headers
    headers = [("SEVERITY", Inches(1.5)),
               ("EXAMPLE", Inches(5.0)),
               ("RESPOND", Inches(2.4)),
               ("RESOLVE", Inches(2.4))]
    hx = rx
    for label, w in headers:
        add_text(s, hx + Inches(0.2), rrow_y, w, Inches(0.3), label,
                 size=9, bold=True, color=ACCENT)
        hx = hx + w
    rrow_y = rrow_y + Inches(0.3)

    for sev, ex, resp, resv, bg, fg in sla:
        add_rect(s, rx, rrow_y, rw, rh, bg,
                 line=(LINE_LIGHT if bg in (BG_MUTED, WHITE) else None))
        cx = rx
        for value, w, bold in [(sev, Inches(1.5), True),
                               (ex, Inches(5.0), False),
                               (resp, Inches(2.4), True),
                               (resv, Inches(2.4), False)]:
            add_text(s, cx + Inches(0.2), rrow_y + Inches(0.06),
                     w - Inches(0.2), Inches(0.3), value,
                     size=11, bold=bold, color=fg)
            cx = cx + w
        rrow_y = rrow_y + rh + Inches(0.04)


# Slide 8 — Security highlights (icon grid)
def slide_security():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Security", "Built in, not bolted on", 8, TOTAL)

    add_text(s, Inches(0.6), Inches(2.6), Inches(12), Inches(0.4),
             "GDPR-aligned by design · the data is sensitive · the controls are real",
             size=14, color=TEXT_MUTED)

    items = [
        ("Encryption",      "TLS in transit · AES-256 at rest"),
        ("MFA",             "Mandatory for all admin accounts"),
        ("RBAC",            "Applicants see only their own data"),
        ("Data Minimisation","Names hidden by default · reveals audit-logged"),
        ("Document Access", "Pre-signed time-limited URLs"),
        ("Audit Trail",     "Immutable log of every CRUD operation"),
        ("UK / EU Hosting", "GDPR-compliant data residency"),
        ("Retention",       "7-year retention · right-to-deletion"),
        ("Backups",         "Daily automated · 30-day PITR"),
        ("OWASP",           "Top-10 hardening · CSP · rate limiting"),
        ("Virus Scanning",  "All uploads scanned before storage"),
        ("Independent Review","GDPR sign-off at G3 before go-live"),
    ]

    cols = 4
    rows = 3
    cell_w = Inches(3.05)
    cell_h = Inches(1.25)
    gap = Inches(0.12)
    start_x = Inches(0.6)
    start_y = Inches(3.3)

    for i, (title, desc) in enumerate(items):
        col = i % cols
        row = i // cols
        x = start_x + (cell_w + gap) * col
        y = start_y + (cell_h + gap) * row
        add_rect(s, x, y, cell_w, cell_h, WHITE, line=LINE_LIGHT)
        # Left accent stripe
        add_rect(s, x, y, Inches(0.06), cell_h, ACCENT)
        add_text(s, x + Inches(0.25), y + Inches(0.2),
                 cell_w - Inches(0.4), Inches(0.4),
                 title, size=12, bold=True, color=NAVY)
        add_text(s, x + Inches(0.25), y + Inches(0.55),
                 cell_w - Inches(0.4), Inches(0.65),
                 desc, size=10, color=TEXT_MUTED, line_spacing=1.25)


# Slide 9 — Migration
def slide_migration():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Migration", "From Grant Tracker · custom workstream", 9, TOTAL)

    add_text(s, Inches(0.6), Inches(2.6), Inches(12), Inches(0.4),
             "Grant Tracker is sunset 31 Dec 2026 · we have runway, but discovery comes first",
             size=14, color=TEXT_MUTED)

    # Three-step process visual
    step_y = Inches(3.3)
    step_h = Inches(2.1)
    step_w = Inches(4.0)
    gap = Inches(0.13)
    start_x = (SW - step_w * 3 - gap * 2) / 2

    steps = [
        ("01", "Discovery Sprint",
         "1–2 weeks · fixed price",
         "Audit Grant Tracker exports.\nDocument data shape.\nIdentify migration risks."),
        ("02", "Fixed Quote",
         "Based on findings",
         "Extract → transform → load.\nScope confirmed.\nTimeline aligned with go-live."),
        ("03", "Cutover",
         "100% migrated",
         "Active accounts moved.\n20% spot-checked.\nReference numbers preserved (WS-/TS-)."),
    ]

    for i, (num, title, sub, desc) in enumerate(steps):
        x = start_x + (step_w + gap) * i
        add_rect(s, x, step_y, step_w, step_h, BG_MUTED)
        add_rect(s, x, step_y, step_w, Inches(0.06), ACCENT)
        add_text(s, x + Inches(0.3), step_y + Inches(0.22),
                 step_w - Inches(0.6), Inches(0.5), num,
                 size=28, bold=True, color=ACCENT, line_spacing=1.0)
        add_text(s, x + Inches(0.3), step_y + Inches(0.78),
                 step_w - Inches(0.6), Inches(0.4), title,
                 size=15, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), step_y + Inches(1.12),
                 step_w - Inches(0.6), Inches(0.3), sub,
                 size=10, italic=True, color=TEXT_MUTED)
        add_text(s, x + Inches(0.3), step_y + Inches(1.45),
                 step_w - Inches(0.6), Inches(0.65), desc,
                 size=10, color=TEXT_DARK, line_spacing=1.3)

    # What needs migrating — bottom strip
    strip_y = Inches(5.7)
    add_text(s, Inches(0.6), strip_y, Inches(12), Inches(0.3),
             "WHAT MOVES ACROSS", size=10, bold=True, color=ACCENT)
    chips = ["Active bursary accounts", "Historical assessments",
             "Documents", "Sibling linkages", "Benchmark fees"]
    cx = Inches(0.6)
    cy = strip_y + Inches(0.32)
    for chip in chips:
        # measure roughly: ~0.13" per char + padding
        w = Inches(0.32 + 0.10 * len(chip))
        add_rect(s, cx, cy, w, Inches(0.42), NAVY)
        add_text(s, cx + Inches(0.22), cy + Inches(0.09),
                 w - Inches(0.4), Inches(0.3), chip,
                 size=11, color=WHITE)
        cx = cx + w + Inches(0.12)


# Slide 10 — Next steps
def slide_next_steps():
    s = prs.slides.add_slide(BLANK)
    # Background
    add_rect(s, 0, 0, SW, SH, NAVY)
    # Right cyan strip
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                SW - Inches(0.18), 0, Inches(0.18), SH)
    accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    # Logo
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.7),
                             Inches(0.36), Inches(0.36))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.fill.background()
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             Inches(1.02), Inches(0.92),
                             Inches(0.10), Inches(0.10))
    dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
    dot.line.fill.background()
    add_text(s, Inches(1.30), Inches(0.74), Inches(4), Inches(0.32),
             "MERIDIAN", size=12, bold=True, color=WHITE)
    add_text(s, Inches(1.30), Inches(0.92), Inches(4), Inches(0.22),
             "TECHNOLOGY GROUP", size=8, color=ACCENT_GLOW)

    # Eyebrow
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.8), Inches(2.0),
                             Inches(0.7), Emu(int(3.5 * 12700)))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    add_text(s, Inches(0.8), Inches(2.1), Inches(8), Inches(0.4),
             "NEXT STEPS", size=11, bold=True, color=ACCENT)

    add_text(s, Inches(0.8), Inches(2.45), Inches(11.5), Inches(0.9),
             "Where we go from here",
             size=42, bold=True, color=WHITE, line_spacing=1.0)

    # Numbered list
    steps = [
        "Confirm acceptance of commercial terms",
        "Schedule UAT kick-off — w/c 18 May 2026",
        "Initiate Grant Tracker migration discovery",
        "Confirm Freshdesk onboarding",
        "Sign go-live date — end of June 2026",
    ]
    sy = Inches(3.7)
    for i, step in enumerate(steps):
        # Number circle
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(0.8), sy + Inches(0.05),
                                  Inches(0.42), Inches(0.42))
        circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT
        circ.line.fill.background()
        add_text(s, Inches(0.8), sy + Inches(0.08), Inches(0.42),
                 Inches(0.36), str(i + 1), size=14, bold=True,
                 color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, Inches(1.45), sy + Inches(0.1), Inches(11), Inches(0.4),
                 step, size=16, color=WHITE)
        sy = sy + Inches(0.5)

    # Footer contact
    add_rect(s, Inches(0.8), Inches(6.7), Inches(11.5), Emu(int(0.5 * 12700)),
             ACCENT)
    add_text(s, Inches(0.8), Inches(6.85), Inches(8), Inches(0.3),
             "Brian Wagner  ·  Meridian Technology Group",
             size=11, bold=True, color=ACCENT)
    add_text(s, SW - Inches(4), Inches(6.85), Inches(3.2), Inches(0.3),
             "meridiantech.group",
             size=11, color=ACCENT_GLOW, align=PP_ALIGN.RIGHT)


# Build
slide_title()
slide_status()
slide_timeline()
slide_commercial_headline()
slide_license_breakdown()
slide_maintenance_scope()
slide_support()
slide_security()
slide_migration()
slide_next_steps()

out = "jwf-bursary-proposal.pptx"
prs.save(out)
print(f"Wrote {out}")
